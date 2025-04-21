#!/usr/bin/env python
# -*- coding: utf-8 -*-

# Copyright 2025 GaoXiaoLon
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.
#
# 创建日期: 2025年4月21日23:17

import os
import argparse
import requests
import json
from tqdm import tqdm
from dotenv import load_dotenv
try:
    from pptx import Presentation
except ImportError:
    print("错误: 找不到python-pptx库。请使用'pip install python-pptx'安装。")
    Presentation = None
from concurrent.futures import ThreadPoolExecutor
from terminology import terminology_manager

# 加载环境变量
load_dotenv()

class DeepSeekTranslator:
    """使用DeepSeek API进行翻译的类"""
    
    def __init__(self, api_key=None, model="deepseek-chat"):
        self.api_key = api_key or os.getenv("DEEPSEEK_API_KEY")
        if not self.api_key:
            raise ValueError("DeepSeek API密钥未提供。请在.env文件中设置DEEPSEEK_API_KEY或作为参数传递。")
        
        self.model = model
        self.api_url = "https://api.deepseek.com/v1/chat/completions"
        self.headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {self.api_key}"
        }
        # 添加翻译内存
        self.translation_memory = {}

    def extract_content(self, text):
        """
        提取文本中的实际内容，去除格式指令和空白内容
        
        Args:
            text (str): 原始文本
            
        Returns:
            str: 提取后的实际内容
        """
        if not text or text.strip() == "":
            return ""
            
        # 移除常见的PPT模板文本指示符
        patterns_to_remove = [
            "点击此处添加文本", "Click to add text", 
            "点击添加", "Click to add",
            "添加标题", "Add title", 
            "添加副标题", "Add subtitle"
        ]
        
        cleaned_text = text
        for pattern in patterns_to_remove:
            if pattern in cleaned_text:
                cleaned_text = cleaned_text.replace(pattern, "")
                
        return cleaned_text.strip()

    def translate(self, text, source_lang="en", target_lang="zh", domain=None):
        """
        使用DeepSeek API翻译文本
        
        Args:
            text (str): 要翻译的文本
            source_lang (str): 源语言代码
            target_lang (str): 目标语言代码 
            domain (str, optional): 专业领域
        
        Returns:
            str: 翻译后的文本
        """
        # 检查是否有实际内容需要翻译
        cleaned_text = self.extract_content(text)
        if not cleaned_text:
            return text
            
        # 检查翻译内存
        cache_key = f"{cleaned_text}_{source_lang}_{target_lang}_{domain}"
        if cache_key in self.translation_memory:
            return self.translation_memory[cache_key]
            
        # 加载领域术语库
        if domain:
            terminology_manager.load_terminology(domain)
            
        # 构建专业提示
        domain_prompt = ""
        if domain:
            domain_prompts = {
                "computer": "你是一位精通计算机科学专业术语的翻译专家，",
                "os": "你是一位精通操作系统专业术语的翻译专家，",
                "general": ""
            }
            domain_prompt = domain_prompts.get(domain, f"你是一位精通{domain}领域专业术语的翻译专家，")
        
        system_prompt = f"{domain_prompt}请将以下{source_lang}文本翻译成{target_lang}，保持专业准确，同时保留原文的格式和标点符号。翻译时，专业术语应使用目标语言中对应的标准术语。文本可能来自幻灯片，请完整翻译所有内容，不要遗漏任何部分。"
        
        # 添加术语库术语到提示中
        terminology_dict = terminology_manager.terminology_dict
        if terminology_dict and len(terminology_dict) > 0:
            terms = []
            for k, v in terminology_dict.items():
                if k in cleaned_text:  # 只添加文本中出现的术语
                    terms.append(f"{k} = {v}")
            
            if terms:
                terminology_prompt = "请特别注意以下专业术语的翻译：\n" + "\n".join(terms)
                system_prompt += "\n\n" + terminology_prompt
        
        data = {
            "model": self.model,
            "messages": [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": cleaned_text}
            ],
            "temperature": 0.3,  # 低温度以确保翻译的确定性
            "max_tokens": 4096
        }
        
        try:
            response = requests.post(
                self.api_url,
                headers=self.headers,
                data=json.dumps(data)
            )
            response.raise_for_status()
            result = response.json()
            translated_text = result["choices"][0]["message"]["content"]
            
            # 使用术语库增强翻译结果
            enhanced_translation = terminology_manager.enhance_translation(cleaned_text, translated_text)
            
            # 保存到翻译内存
            self.translation_memory[cache_key] = enhanced_translation
            
            return enhanced_translation
        except Exception as e:
            print(f"翻译时出错: {e}")
            return text  # 出错时返回原文
            
    def batch_translate(self, texts, source_lang="en", target_lang="zh", domain=None):
        """
        批量翻译多段文本
        
        Args:
            texts (list): 要翻译的文本列表
            source_lang (str): 源语言代码
            target_lang (str): 目标语言代码 
            domain (str, optional): 专业领域
        
        Returns:
            list: 翻译后的文本列表
        """
        if not texts:
            return []
            
        # 过滤出需要翻译的非空文本
        filtered_texts = []
        original_indices = []
        
        for i, text in enumerate(texts):
            cleaned_text = self.extract_content(text)
            if cleaned_text:
                filtered_texts.append(cleaned_text)
                original_indices.append(i)
                
        if not filtered_texts:
            return texts
            
        # 合并所有文本，用特殊标记分隔
        separator = "\n===[SEPARATOR]===\n"
        merged_text = separator.join(filtered_texts)
        
        # 构建专业提示
        domain_prompt = ""
        if domain:
            terminology_manager.load_terminology(domain)
            domain_prompts = {
                "computer": "你是一位精通计算机科学专业术语的翻译专家，",
                "os": "你是一位精通操作系统专业术语的翻译专家，",
                "general": ""
            }
            domain_prompt = domain_prompts.get(domain, f"你是一位精通{domain}领域专业术语的翻译专家，")
        
        system_prompt = f"{domain_prompt}以下是多个由特殊标记'===[SEPARATOR]==='分隔的{source_lang}文本段落。请将每个段落翻译成{target_lang}，并在翻译文本之间保留相同的分隔标记。保持专业准确，同时保留原文的格式和标点符号。翻译时，专业术语应使用目标语言中对应的标准术语。确保完整翻译所有内容，不要遗漏任何部分。"
        
        # 添加术语库术语到提示中
        terminology_dict = terminology_manager.terminology_dict
        if terminology_dict and len(terminology_dict) > 0:
            terms = []
            for k, v in terminology_dict.items():
                if k in merged_text:  # 只添加文本中出现的术语
                    terms.append(f"{k} = {v}")
            
            if terms:
                terminology_prompt = "请特别注意以下专业术语的翻译：\n" + "\n".join(terms)
                system_prompt += "\n\n" + terminology_prompt
        
        data = {
            "model": self.model,
            "messages": [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": merged_text}
            ],
            "temperature": 0.3,
            "max_tokens": 8192  # 增加token上限以适应更多内容
        }
        
        try:
            response = requests.post(
                self.api_url,
                headers=self.headers,
                data=json.dumps(data)
            )
            response.raise_for_status()
            result = response.json()
            translated_merged = result["choices"][0]["message"]["content"]
            
            # 使用术语库增强翻译结果
            enhanced_translation = terminology_manager.enhance_translation(merged_text, translated_merged)
            
            # 分割翻译后的文本
            translated_texts = enhanced_translation.split(separator)
            
            # 确保翻译结果与原文数量匹配
            if len(translated_texts) != len(filtered_texts):
                print(f"警告: 翻译段落数量不匹配 (原文: {len(filtered_texts)}, 译文: {len(translated_texts)})")
                # 如果段落数量不匹配，则单独翻译每个段落
                translated_texts = [self.translate(text, source_lang, target_lang, domain) for text in filtered_texts]
            
            # 将翻译结果插回原始位置
            result_texts = texts.copy()
            for i, idx in enumerate(original_indices):
                if i < len(translated_texts):
                    result_texts[idx] = translated_texts[i]
                    
                    # 保存到翻译内存
                    cache_key = f"{filtered_texts[i]}_{source_lang}_{target_lang}_{domain}"
                    self.translation_memory[cache_key] = translated_texts[i]
            
            return result_texts
            
        except Exception as e:
            print(f"批量翻译时出错: {e}")
            # 回退到单独翻译
            result_texts = texts.copy()
            for i, idx in enumerate(original_indices):
                result_texts[idx] = self.translate(filtered_texts[i], source_lang, target_lang, domain)
            return result_texts

class PPTTranslator:
    """PPT文件翻译器类"""
    
    def __init__(self, translator, source_lang="en", target_lang="zh", domain=None):
        if Presentation is None:
            raise ImportError("请安装python-pptx库以使用此功能")
        
        self.translator = translator
        self.source_lang = source_lang
        self.target_lang = target_lang
        self.domain = domain
        
    def translate_ppt(self, input_file, output_file, progress_callback=None):
        """
        翻译PPT文件
        
        Args:
            input_file (str): 输入PPT文件路径
            output_file (str): 输出PPT文件路径
            progress_callback (function, optional): 进度回调函数
        """
        # 加载PPT
        prs = Presentation(input_file)
        total_slides = len(prs.slides)
        
        print(f"开始翻译PPT: {input_file}")
        print(f"总计 {total_slides} 张幻灯片")
        
        # 遍历所有幻灯片和形状进行翻译
        for i, slide in enumerate(tqdm(prs.slides, desc="翻译幻灯片")):
            # 更新进度
            if progress_callback:
                progress_callback(i+1, total_slides)
            
            # 递归处理每个幻灯片上的所有形状（包括嵌套形状）    
            self.process_shapes(slide.shapes, slide_index=i)
        
        # 保存翻译后的PPT
        prs.save(output_file)
        print(f"翻译完成，已保存至: {output_file}")
    
    def process_shapes(self, shapes, slide_index=0):
        """
        递归处理形状集合中的所有形状
        
        Args:
            shapes: 形状集合
            slide_index: 当前幻灯片索引（用于错误报告）
        """
        for shape in shapes:
            try:
                # 处理常规文本框
                if hasattr(shape, "text_frame") and shape.text_frame:
                    self.translate_text_frame(shape.text_frame)
                
                # 处理表格文本
                if hasattr(shape, "table"):
                    self.translate_table(shape.table)
                
                # 处理图表中的文本 (标题、轴标签等)
                if hasattr(shape, "chart"):
                    self.translate_chart(shape, slide_index)
                
                # 处理SmartArt (尝试通过访问文本属性)
                if hasattr(shape, "shape_type") and shape.shape_type == 14:  # MSO_SHAPE_TYPE.SMART_ART
                    self.translate_smart_art(shape, slide_index)
                
                # 处理WordArt (通常作为特殊的文本框实现)
                if hasattr(shape, "shape_type") and shape.shape_type == 7:  # MSO_SHAPE_TYPE.TEXT_BOX
                    self.translate_word_art(shape, slide_index)
                
                # 处理组合形状 - 关键的增强部分
                if hasattr(shape, "group_items") and shape.group_items:
                    self.process_shapes(shape.group_items, slide_index)
                
                # 处理占位符中的形状
                if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                    if hasattr(shape, "placeholder_format") and hasattr(shape.placeholder_format, "idx"):
                        try:
                            placeholder = shape.placeholder_format
                            if hasattr(placeholder, "type") and placeholder.type:
                                print(f"处理占位符 (幻灯片 {slide_index+1}): 类型 {placeholder.type}")
                        except Exception as e:
                            print(f"无法识别占位符类型 (幻灯片 {slide_index+1}): {e}")
            
            except Exception as e:
                print(f"处理形状时出错 (幻灯片 {slide_index+1}): {e}")
    
    def translate_text_frame(self, text_frame):
        """
        翻译文本框内容
        
        Args:
            text_frame: 文本框对象
        """
        if not text_frame.text.strip():
            return
        
        # 收集所有段落文本
        paragraphs = []
        for paragraph in text_frame.paragraphs:
            para_text = paragraph.text
            if para_text.strip():
                paragraphs.append(para_text)
        
        # 如果没有内容，则直接返回
        if not paragraphs:
            return
            
        # 使用批量翻译API进行翻译
        translated_paragraphs = self.translator.batch_translate(
            paragraphs, 
            self.source_lang, 
            self.target_lang,
            self.domain
        )
        
        # 更新文本框中的段落
        for j, paragraph in enumerate(text_frame.paragraphs):
            if j < len(translated_paragraphs) and paragraph.text.strip():
                # 保留原始格式处理
                if len(paragraph.runs) == 0:
                    # 没有runs，直接添加新的run
                    run = paragraph.add_run()
                    run.text = translated_paragraphs[j]
                else:
                    # 更新现有runs，保留格式
                    for run_idx, run in enumerate(paragraph.runs):
                        if run_idx == 0:
                            run.text = translated_paragraphs[j]
                        else:
                            run.text = ""
    
    def translate_table(self, table):
        """
        翻译表格内容
        
        Args:
            table: 表格对象
        """
        # 收集所有单元格的文本
        cell_texts = []
        cell_mapping = []  # 存储(row, col)到文本列表索引的映射
        
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                text = cell.text.strip()
                if text:
                    cell_texts.append(text)
                    cell_mapping.append((row_idx, col_idx))
        
        # 如果没有内容，则直接返回
        if not cell_texts:
            return
            
        # 使用批量翻译
        translated_texts = self.translator.batch_translate(
            cell_texts, 
            self.source_lang, 
            self.target_lang,
            self.domain
        )
        
        # 更新表格单元格文本
        for idx, (row_idx, col_idx) in enumerate(cell_mapping):
            if idx < len(translated_texts):
                cell = table.rows[row_idx].cells[col_idx]
                # 清除原始内容
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = ""
                
                # 填入翻译后的文本
                first_paragraph = cell.text_frame.paragraphs[0]
                if not first_paragraph.runs:
                    run = first_paragraph.add_run()
                    run.text = translated_texts[idx]
                else:
                    first_paragraph.runs[0].text = translated_texts[idx]
    
    def translate_chart(self, shape, slide_index):
        """
        翻译图表中的文字内容，包括标题、坐标轴标签、图例和数据标签
        
        Args:
            shape: Chart对象
            slide_index: 当前幻灯片索引（用于错误报告）
        """
        try:
            if not hasattr(shape, "chart"):
                return
                
            chart = shape.chart
            texts_to_translate = []
            text_elements = []
            
            # 收集图表标题
            if hasattr(chart, "chart_title") and chart.chart_title and hasattr(chart.chart_title, "text_frame"):
                if chart.chart_title.text_frame.text.strip():
                    texts_to_translate.append(chart.chart_title.text_frame.text)
                    text_elements.append(chart.chart_title.text_frame)
            
            # 收集X轴标题
            if hasattr(chart, "has_category_axis") and chart.has_category_axis:
                category_axis = chart.category_axis
                if hasattr(category_axis, "axis_title") and category_axis.axis_title:
                    if hasattr(category_axis.axis_title, "text_frame") and category_axis.axis_title.text_frame.text.strip():
                        texts_to_translate.append(category_axis.axis_title.text_frame.text)
                        text_elements.append(category_axis.axis_title.text_frame)
            
            # 收集Y轴标题
            if hasattr(chart, "has_value_axis") and chart.has_value_axis:
                value_axis = chart.value_axis
                if hasattr(value_axis, "axis_title") and value_axis.axis_title:
                    if hasattr(value_axis.axis_title, "text_frame") and value_axis.axis_title.text_frame.text.strip():
                        texts_to_translate.append(value_axis.axis_title.text_frame.text)
                        text_elements.append(value_axis.axis_title.text_frame)
            
            # 收集图例中的系列名称
            if hasattr(chart, "series"):
                for series in chart.series:
                    if hasattr(series, "name") and series.name and series.name.strip():
                        texts_to_translate.append(series.name)
                        text_elements.append(("series_name", series))
            
            # 尝试收集类别标签
            if hasattr(chart, "plots") and chart.plots:
                for plot in chart.plots:
                    if hasattr(plot, "categories"):
                        for i, category in enumerate(plot.categories):
                            if category and category.strip():
                                texts_to_translate.append(category)
                                text_elements.append(("category", (plot, i)))
            
            # 尝试收集数据标签
            if hasattr(chart, "plots") and chart.plots:
                for plot in chart.plots:
                    if hasattr(plot, "series"):
                        for series in plot.series:
                            if hasattr(series, "data_labels"):
                                for i, label in enumerate(series.data_labels):
                                    if hasattr(label, "text_frame") and label.text_frame.text.strip():
                                        texts_to_translate.append(label.text_frame.text)
                                        text_elements.append(label.text_frame)
            
            # 如果没有要翻译的文本，直接返回
            if not texts_to_translate:
                return
                
            # 批量翻译收集到的文本
            translated_texts = self.translator.batch_translate(
                texts_to_translate,
                self.source_lang,
                self.target_lang,
                self.domain
            )
            
            # 更新文本
            if translated_texts and len(translated_texts) == len(text_elements):
                for i, element in enumerate(text_elements):
                    if isinstance(element, tuple):
                        element_type, element_data = element
                        if element_type == "series_name":
                            element_data.name = translated_texts[i]
                        elif element_type == "category":
                            plot, category_index = element_data
                            if hasattr(plot, "categories") and category_index < len(plot.categories):
                                plot.categories[category_index] = translated_texts[i]
                    else:
                        element.text = translated_texts[i]
                
        except Exception as e:
            print(f"无法翻译图表 (幻灯片 {slide_index+1}): {e}")
    
    def translate_smart_art(self, shape, slide_index):
        """
        翻译SmartArt内容
        
        Args:
            shape: SmartArt形状对象
            slide_index: 当前幻灯片索引（用于错误报告）
        """
        try:
            # SmartArt可能包含多个文本框或嵌套形状
            # 查找所有可能包含文本的元素
            texts = []
            text_elements = []
            
            # 尝试从SmartArt的XML结构中提取文本
            try:
                # 遍历SmartArt元素树
                for node in shape.element.iter():
                    # 检查节点是否包含文本属性
                    if hasattr(node, 'text') and node.text and node.text.strip():
                        texts.append(node.text)
                        text_elements.append(node)
            except:
                pass
                
            # 如果上面的方法没有找到文本，尝试其他方法
            if not texts:
                # 尝试查找组合形状内的文本框
                if hasattr(shape, 'shapes'):
                    for inner_shape in shape.shapes:
                        if hasattr(inner_shape, 'text_frame') and inner_shape.text_frame:
                            if inner_shape.text_frame.text.strip():
                                texts.append(inner_shape.text_frame.text)
                                text_elements.append(inner_shape.text_frame)
                
                # 尝试从图形数据中获取文本
                if hasattr(shape, 'data'):
                    data = shape.data
                    if hasattr(data, 'text') and data.text.strip():
                        texts.append(data.text)
                        text_elements.append(data)
            
            # 如果没有可翻译的文本，直接返回
            if not texts:
                return
            
            # 批量翻译所有收集到的文本
            translated_texts = self.translator.batch_translate(
                texts,
                self.source_lang,
                self.target_lang,
                self.domain
            )
            
            # 更新文本元素
            if translated_texts and len(translated_texts) == len(text_elements):
                for i, element in enumerate(text_elements):
                    # 根据元素类型设置翻译后的文本
                    if hasattr(element, 'text'):
                        element.text = translated_texts[i]
                    elif hasattr(element, 'text_frame'):
                        element.text = translated_texts[i]
                        
        except Exception as e:
            print(f"无法翻译SmartArt (幻灯片 {slide_index+1}): {e}")
    
    def translate_word_art(self, shape, slide_index):
        """
        翻译WordArt内容（艺术字）及其他特殊文本形状
        
        Args:
            shape: WordArt形状对象或包含特殊文本的形状
            slide_index: 当前幻灯片索引（用于错误报告）
        """
        try:
            texts_to_translate = []
            text_properties = []
            
            # 检查不同类型的文本属性
            if hasattr(shape, "text") and shape.text.strip():
                texts_to_translate.append(shape.text.strip())
                text_properties.append(("text", None))
                
            # 处理文本框架
            elif hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text") and shape.text_frame.text.strip():
                texts_to_translate.append(shape.text_frame.text.strip())
                text_properties.append(("text_frame", None))
                
            # 处理段落
            elif hasattr(shape, "text_frame") and hasattr(shape.text_frame, "paragraphs"):
                for i, para in enumerate(shape.text_frame.paragraphs):
                    if para.text.strip():
                        texts_to_translate.append(para.text.strip())
                        text_properties.append(("paragraph", i))
                
            # 如果没有要翻译的文本，直接返回
            if not texts_to_translate:
                return
            
            # 批量翻译所有收集到的文本
            translated_texts = self.translator.batch_translate(
                texts_to_translate,
                self.source_lang,
                self.target_lang,
                self.domain
            )
            
            # 更新各种文本属性
            if translated_texts and len(translated_texts) > 0:
                for i, (prop_type, index) in enumerate(text_properties):
                    if i < len(translated_texts):
                        if prop_type == "text":
                            shape.text = translated_texts[i]
                        elif prop_type == "text_frame":
                            shape.text_frame.text = translated_texts[i]
                        elif prop_type == "paragraph" and index is not None:
                            shape.text_frame.paragraphs[index].text = translated_texts[i]
                            
        except Exception as e:
            print(f"无法翻译WordArt (幻灯片 {slide_index+1}): {e}")

def main():
    """主函数"""
    parser = argparse.ArgumentParser(description="PPT翻译工具 - 使用DeepSeek API进行专业翻译")
    parser.add_argument("--input", required=True, help="输入PPT文件路径")
    parser.add_argument("--output", required=True, help="输出PPT文件路径")
    parser.add_argument("--source", default="en", help="源语言代码，默认为'en'")
    parser.add_argument("--target", default="zh", help="目标语言代码，默认为'zh'")
    parser.add_argument("--api-key", help="DeepSeek API密钥（可选，也可在.env文件中设置）")
    parser.add_argument("--model", default="deepseek-chat", help="DeepSeek模型名称")
    parser.add_argument("--domain", choices=["computer", "os", "general"], 
                        help="专业领域，可选值: computer (计算机), os (操作系统), general (通用)")
    
    args = parser.parse_args()
    
    # 检查输入文件
    if not os.path.exists(args.input):
        print(f"错误: 输入文件不存在: {args.input}")
        return
    
    try:
        # 初始化翻译器
        translator = DeepSeekTranslator(api_key=args.api_key, model=args.model)
        ppt_translator = PPTTranslator(
            translator, 
            source_lang=args.source, 
            target_lang=args.target,
            domain=args.domain
        )
        
        # 翻译PPT
        ppt_translator.translate_ppt(args.input, args.output)
        
    except ImportError as e:
        print(f"导入错误: {e}")
        print("请安装必要的依赖: pip install python-pptx")
    except Exception as e:
        print(f"翻译过程中出错: {e}")

if __name__ == "__main__":
    main() 
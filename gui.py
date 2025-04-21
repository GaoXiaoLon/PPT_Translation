#!/usr/bin/env python
# -*- coding: utf-8 -*-

import os
import sys
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from dotenv import load_dotenv
import threading
import time

# 检查必要的依赖
try:
    from pptx import Presentation
    from translator import DeepSeekTranslator, PPTTranslator
except ImportError as e:
    def show_error_and_exit():
        root = tk.Tk()
        root.withdraw()
        messagebox.showerror(
            "缺少依赖", 
            f"错误: {str(e)}\n\n请运行'安装依赖.bat'或使用以下命令安装:\npip install python-pptx requests python-dotenv tqdm"
        )
        sys.exit(1)
    show_error_and_exit()

# 加载环境变量
load_dotenv()

class PPTTranslatorApp:
    """PPT翻译工具的图形用户界面"""
    
    def __init__(self, root):
        self.root = root
        self.root.title("PPT 专业翻译工具 - DeepSeek版")
        self.root.geometry("800x640")  # 更大的窗口尺寸
        self.root.resizable(True, True)
        self.root.configure(bg="#f5f5f5")  # 设置背景色
        
        # 设置样式
        self.style = ttk.Style()
        self.style.theme_use('clam')  # 使用更现代的主题
        self.style.configure("TFrame", background="#f5f5f5")
        self.style.configure("TLabelframe", background="#f5f5f5")
        self.style.configure("TLabelframe.Label", background="#f5f5f5", font=('Arial', 11, 'bold'))
        self.style.configure("TLabel", background="#f5f5f5", font=('Arial', 10))
        self.style.configure("TButton", padding=6, font=('Arial', 10))
        self.style.configure("Header.TLabel", font=('Arial', 16, 'bold'), background="#f5f5f5")
        self.style.configure("BigButton.TButton", padding=12, font=('Arial', 13, 'bold'))
        self.style.configure("TopButton.TButton", padding=10, font=('Arial', 12, 'bold'), background="#4CAF50")
        self.style.configure("Green.TButton", background="#4CAF50", foreground="white")
        
        # 创建主框架
        main_frame = ttk.Frame(root)
        main_frame.pack(fill=tk.BOTH, expand=True, padx=25, pady=20)
        
        # 标题区域和大翻译按钮一起放在顶部框架中
        top_frame = ttk.Frame(main_frame)
        top_frame.pack(fill=tk.X, pady=(0, 15))
        
        header_label = ttk.Label(
            top_frame, 
            text="PPT 专业翻译工具", 
            style="Header.TLabel"
        )
        header_label.pack(side=tk.LEFT)
        
        # 大的开始翻译按钮放在标题右侧
        self.translate_big_btn = ttk.Button(
            top_frame, 
            text="开始翻译",
            style="TopButton.TButton",
            command=self.start_translation,
            width=15
        )
        self.translate_big_btn.pack(side=tk.RIGHT, padx=10, pady=0)
        
        # API设置框架
        api_frame = ttk.LabelFrame(main_frame, text="DeepSeek API 设置")
        api_frame.pack(fill=tk.X, pady=10)
        
        # 创建网格布局内的内部框架，便于添加间距
        api_inner_frame = ttk.Frame(api_frame)
        api_inner_frame.pack(fill=tk.X, padx=15, pady=10)
        
        ttk.Label(api_inner_frame, text="API 密钥:").grid(row=0, column=0, sticky=tk.W, pady=8)
        
        # 从环境变量获取API密钥
        default_api_key = os.getenv("DEEPSEEK_API_KEY", "")
        
        self.api_key_var = tk.StringVar(value=default_api_key)
        self.api_key_entry = ttk.Entry(api_inner_frame, textvariable=self.api_key_var, width=55, show="*")
        self.api_key_entry.grid(row=0, column=1, sticky=tk.W, padx=5, pady=8)
        
        # 显示/隐藏密钥按钮
        self.show_key = tk.BooleanVar(value=False)
        self.show_key_btn = ttk.Checkbutton(
            api_inner_frame, 
            text="显示", 
            variable=self.show_key, 
            command=self.toggle_key_visibility
        )
        self.show_key_btn.grid(row=0, column=2, sticky=tk.W, padx=5, pady=8)
        
        # 模型选择
        ttk.Label(api_inner_frame, text="模型:").grid(row=1, column=0, sticky=tk.W, pady=8)
        self.model_var = tk.StringVar(value="deepseek-chat")
        self.model_combo = ttk.Combobox(api_inner_frame, textvariable=self.model_var, width=25, state="readonly")
        self.model_combo['values'] = ('deepseek-chat', 'deepseek-coder')
        self.model_combo.grid(row=1, column=1, sticky=tk.W, padx=5, pady=8)
        
        # 翻译设置框架
        translation_frame = ttk.LabelFrame(main_frame, text="翻译设置")
        translation_frame.pack(fill=tk.X, pady=15)
        
        # 创建内部框架
        translation_inner_frame = ttk.Frame(translation_frame)
        translation_inner_frame.pack(fill=tk.X, padx=15, pady=10)
        
        # 源语言
        ttk.Label(translation_inner_frame, text="源语言:").grid(row=0, column=0, sticky=tk.W, pady=8)
        self.source_lang_var = tk.StringVar(value="en")
        self.source_lang_combo = ttk.Combobox(translation_inner_frame, textvariable=self.source_lang_var, width=12, state="readonly")
        self.source_lang_combo['values'] = ('en', 'zh', 'ja', 'ko', 'fr', 'de', 'es', 'ru')
        self.source_lang_combo.grid(row=0, column=1, sticky=tk.W, padx=5, pady=8)
        
        # 目标语言
        ttk.Label(translation_inner_frame, text="目标语言:").grid(row=0, column=2, sticky=tk.W, pady=8, padx=(20, 0))
        self.target_lang_var = tk.StringVar(value="zh")
        self.target_lang_combo = ttk.Combobox(translation_inner_frame, textvariable=self.target_lang_var, width=12, state="readonly")
        self.target_lang_combo['values'] = ('zh', 'en', 'ja', 'ko', 'fr', 'de', 'es', 'ru')
        self.target_lang_combo.grid(row=0, column=3, sticky=tk.W, padx=5, pady=8)
        
        # 专业领域
        ttk.Label(translation_inner_frame, text="专业领域:").grid(row=1, column=0, sticky=tk.W, pady=8)
        self.domain_var = tk.StringVar(value="general")
        self.domain_combo = ttk.Combobox(translation_inner_frame, textvariable=self.domain_var, width=18, state="readonly")
        self.domain_combo['values'] = ('general', 'computer', 'os', 'medicine', 'law', 'finance')
        self.domain_combo.grid(row=1, column=1, sticky=tk.W, padx=5, pady=8, columnspan=2)
        
        # 文件选择框架
        file_frame = ttk.LabelFrame(main_frame, text="文件选择")
        file_frame.pack(fill=tk.X, pady=15)
        
        # 创建内部框架
        file_inner_frame = ttk.Frame(file_frame)
        file_inner_frame.pack(fill=tk.X, padx=15, pady=10)
        
        # 输入文件
        ttk.Label(file_inner_frame, text="输入 PPT:").grid(row=0, column=0, sticky=tk.W, pady=8)
        self.input_file_var = tk.StringVar()
        self.input_file_entry = ttk.Entry(file_inner_frame, textvariable=self.input_file_var, width=55)
        self.input_file_entry.grid(row=0, column=1, sticky=tk.W, padx=5, pady=8)
        
        self.browse_input_btn = ttk.Button(
            file_inner_frame, 
            text="浏览...", 
            command=self.browse_input_file
        )
        self.browse_input_btn.grid(row=0, column=2, sticky=tk.W, padx=5, pady=8)
        
        # 输出文件
        ttk.Label(file_inner_frame, text="输出 PPT:").grid(row=1, column=0, sticky=tk.W, pady=8)
        self.output_file_var = tk.StringVar()
        self.output_file_entry = ttk.Entry(file_inner_frame, textvariable=self.output_file_var, width=55)
        self.output_file_entry.grid(row=1, column=1, sticky=tk.W, padx=5, pady=8)
        
        self.browse_output_btn = ttk.Button(
            file_inner_frame, 
            text="浏览...", 
            command=self.browse_output_file
        )
        self.browse_output_btn.grid(row=1, column=2, sticky=tk.W, padx=5, pady=8)
        
        # 翻译进度框架
        progress_frame = ttk.LabelFrame(main_frame, text="翻译进度")
        progress_frame.pack(fill=tk.X, pady=15)
        
        # 创建内部框架
        progress_inner_frame = ttk.Frame(progress_frame)
        progress_inner_frame.pack(fill=tk.X, padx=15, pady=10)
        
        # 状态标签
        self.status_var = tk.StringVar(value="就绪")
        status_label = ttk.Label(
            progress_inner_frame, 
            textvariable=self.status_var, 
            font=('Arial', 11, 'bold')
        )
        status_label.pack(fill=tk.X, pady=(0, 10))
        
        # 翻译进度条
        self.progress_var = tk.DoubleVar()
        self.progress_bar = ttk.Progressbar(
            progress_inner_frame, 
            variable=self.progress_var,
            maximum=100,
            length=100,
            mode='determinate',
            style="TProgressbar"
        )
        self.progress_bar.pack(fill=tk.X, pady=5)
        
        # 日志框架
        log_frame = ttk.LabelFrame(main_frame, text="翻译日志")
        log_frame.pack(fill=tk.BOTH, expand=True, pady=15)
        
        # 创建内部框架
        log_inner_frame = ttk.Frame(log_frame)
        log_inner_frame.pack(fill=tk.BOTH, expand=True, padx=15, pady=10)
        
        # 滚动条
        scrollbar = ttk.Scrollbar(log_inner_frame)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        
        # 日志文本框
        self.log_text = tk.Text(
            log_inner_frame, 
            height=8, 
            yscrollcommand=scrollbar.set,
            font=('Consolas', 9),
            bg='#f9f9f9',
            wrap=tk.WORD
        )
        self.log_text.pack(fill=tk.BOTH, expand=True)
        scrollbar.config(command=self.log_text.yview)
        
        # 设置样式
        self.style.configure("TProgressbar", thickness=25)
        
        # 初始状态
        self.translation_in_progress = False
        self.log("PPT专业翻译工具已启动，请配置参数并选择文件。")
    
    def toggle_key_visibility(self):
        """切换API密钥的可见性"""
        if self.show_key.get():
            self.api_key_entry.config(show="")
        else:
            self.api_key_entry.config(show="*")
    
    def browse_input_file(self):
        """浏览并选择输入PPT文件"""
        filename = filedialog.askopenfilename(
            title="选择PPT文件",
            filetypes=(("PowerPoint files", "*.pptx *.ppt"), ("All files", "*.*"))
        )
        if filename:
            self.input_file_var.set(filename)
            # 自动生成输出文件路径
            base_name = os.path.basename(filename)
            name, ext = os.path.splitext(base_name)
            output_name = f"{name}_translated{ext}"
            output_dir = os.path.dirname(filename)
            self.output_file_var.set(os.path.join(output_dir, output_name))
            self.log(f"已选择输入文件: {filename}")
    
    def browse_output_file(self):
        """浏览并选择输出PPT文件保存位置"""
        filename = filedialog.asksaveasfilename(
            title="保存翻译后的PPT",
            filetypes=(("PowerPoint files", "*.pptx"), ("All files", "*.*")),
            defaultextension=".pptx"
        )
        if filename:
            self.output_file_var.set(filename)
            self.log(f"已设置输出文件: {filename}")
    
    def log(self, message):
        """向日志添加消息"""
        self.log_text.insert(tk.END, f"[{time.strftime('%H:%M:%S')}] {message}\n")
        self.log_text.see(tk.END)  # 自动滚动到底部
        
    def validate_inputs(self):
        """验证用户输入"""
        if not self.api_key_var.get():
            messagebox.showerror("错误", "请输入DeepSeek API密钥")
            return False
            
        if not self.input_file_var.get():
            messagebox.showerror("错误", "请选择输入PPT文件")
            return False
            
        if not self.output_file_var.get():
            messagebox.showerror("错误", "请指定输出PPT文件")
            return False
            
        if not os.path.exists(self.input_file_var.get()):
            messagebox.showerror("错误", f"输入文件不存在: {self.input_file_var.get()}")
            return False
            
        return True
    
    def start_translation(self):
        """开始翻译过程"""
        if self.translation_in_progress:
            messagebox.showinfo("提示", "翻译任务正在进行中")
            return
            
        if not self.validate_inputs():
            return
            
        # 禁用按钮，防止重复操作
        self.translate_big_btn.config(state=tk.DISABLED)
        self.translation_in_progress = True
        self.status_var.set("翻译中...")
        
        # 创建线程执行翻译任务
        translation_thread = threading.Thread(target=self.translation_task)
        translation_thread.daemon = True
        translation_thread.start()
        
        # 更新按钮颜色
        self.style.configure("TopButton.TButton", background="#ffcc00")
    
    def translation_task(self):
        """在后台线程中执行翻译任务"""
        try:
            self.log(f"开始翻译任务...")
            self.log(f"输入文件: {self.input_file_var.get()}")
            self.log(f"输出文件: {self.output_file_var.get()}")
            self.log(f"源语言: {self.source_lang_var.get()}, 目标语言: {self.target_lang_var.get()}")
            self.log(f"专业领域: {self.domain_var.get()}")
            
            # 初始化翻译器
            self.log("初始化DeepSeek翻译器...")
            translator = DeepSeekTranslator(
                api_key=self.api_key_var.get(), 
                model=self.model_var.get()
            )
            
            # 自定义进度更新回调
            def update_progress(current, total):
                progress = (current / total) * 100
                self.progress_var.set(progress)
                self.status_var.set(f"翻译中... {progress:.1f}%")
                self.log(f"翻译进度: {current}/{total} 张幻灯片 ({progress:.1f}%)")
                # 让界面更新
                self.root.update_idletasks()
                
            # 修改PPTTranslator类处理进度回调
            class CustomPPTTranslator(PPTTranslator):
                def translate_ppt(self, input_file, output_file, progress_callback=None):
                    # 加载PPT
                    prs = Presentation(input_file)
                    total_slides = len(prs.slides)
                    
                    # 遍历所有幻灯片和形状进行翻译
                    for i, slide in enumerate(prs.slides):
                        # 更新进度
                        if progress_callback:
                            progress_callback(i+1, total_slides)
                            
                        for shape in slide.shapes:
                            if hasattr(shape, "text_frame") and shape.text_frame:
                                if shape.text_frame.text.strip():
                                    # 收集所有段落文本
                                    paragraphs = []
                                    for paragraph in shape.text_frame.paragraphs:
                                        para_text = paragraph.text
                                        if para_text.strip():
                                            paragraphs.append(para_text)
                                    
                                    # 批量翻译所有段落文本
                                    if paragraphs:
                                        text_to_translate = "\n".join(paragraphs)
                                        translated_text = self.translator.translate(
                                            text_to_translate, 
                                            self.source_lang, 
                                            self.target_lang,
                                            self.domain
                                        )
                                        
                                        # 分割翻译结果并更新段落
                                        translated_paragraphs = translated_text.split('\n')
                                        for j, paragraph in enumerate(shape.text_frame.paragraphs):
                                            if j < len(translated_paragraphs) and paragraph.text.strip():
                                                # 保留原始格式
                                                for run_idx, run in enumerate(paragraph.runs):
                                                    if run_idx == 0 and translated_paragraphs[j].strip():
                                                        run.text = translated_paragraphs[j]
                                                    elif run_idx > 0:
                                                        run.text = ""
                            
                            # 处理表格文本
                            if hasattr(shape, "table"):
                                for row in shape.table.rows:
                                    for cell in row.cells:
                                        if cell.text.strip():
                                            cell.text = self.translator.translate(
                                                cell.text, 
                                                self.source_lang, 
                                                self.target_lang,
                                                self.domain
                                            )
                    
                    # 保存翻译后的PPT
                    prs.save(output_file)
            
            # 初始化PPT翻译器
            ppt_translator = CustomPPTTranslator(
                translator, 
                source_lang=self.source_lang_var.get(), 
                target_lang=self.target_lang_var.get(),
                domain=self.domain_var.get()
            )
            
            # 执行翻译
            self.log("开始翻译PPT内容...")
            ppt_translator.translate_ppt(
                self.input_file_var.get(), 
                self.output_file_var.get(),
                progress_callback=update_progress
            )
            
            self.status_var.set("翻译完成!")
            self.log(f"翻译完成！已保存至: {self.output_file_var.get()}")
            
            # 设置按钮颜色为绿色
            self.style.configure("TopButton.TButton", background="#4CAF50")
            
            # 显示完成消息框
            messagebox.showinfo("成功", f"PPT翻译已完成！\n文件已保存至: {self.output_file_var.get()}")
            
        except ImportError as e:
            self.log(f"导入错误: {str(e)}")
            self.status_var.set("依赖缺失")
            messagebox.showerror("依赖错误", f"缺少必要的依赖: {str(e)}\n\n请运行'安装依赖.bat'或使用命令:\npip install python-pptx requests python-dotenv tqdm")
            
        except Exception as e:
            self.log(f"翻译过程中出错: {str(e)}")
            self.status_var.set("翻译失败")
            # 设置按钮颜色为红色
            self.style.configure("TopButton.TButton", background="#f44336")
            messagebox.showerror("错误", f"翻译过程中出错: {str(e)}")
            
        finally:
            # 恢复按钮状态
            self.translate_big_btn.config(state=tk.NORMAL)
            self.translation_in_progress = False

def main():
    """主函数"""
    try:
        root = tk.Tk()
        app = PPTTranslatorApp(root)
        root.mainloop()
    except Exception as e:
        # 捕获任何未处理的异常
        messagebox.showerror("错误", f"发生未预期的错误: {str(e)}")
        sys.exit(1)

if __name__ == "__main__":
    main() 